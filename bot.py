import logging
import os
import time
import base64
import requests
import chardet
from datetime import date
import PyPDF2
from telegram.ext import Updater, CommandHandler, MessageHandler, Filters, CallbackContext
from telegram import Update
from bs4 import BeautifulSoup, NavigableString
from googletrans import Translator
import arabic_reshaper
from bidi.algorithm import get_display

# Ø§Ù„Ù…ÙƒØªØ¨Ø§Øª Ø§Ù„Ø®Ø§ØµØ© Ø¨Ù…Ù„ÙØ§Øª DOCX Ùˆ PPTX
from docx import Document
from pptx import Presentation

# Ø¥Ø¹Ø¯Ø§Ø¯ ØªØ³Ø¬ÙŠÙ„ Ø§Ù„Ø£Ø®Ø·Ø§Ø¡
logging.basicConfig(format='%(asctime)s - %(name)s - %(levelname)s - %(message)s', level=logging.INFO)
logger = logging.getLogger(__name__)

# Ø¥Ø¹Ø¯Ø§Ø¯ Ø§Ù„ØªÙˆÙƒÙ† ÙˆÙ…ÙØ§ØªÙŠØ­ API
TELEGRAM_TOKEN = '8060810536:AAFGPiwBQuYSJG0UUwiypPfowr10qqc0nq0'
CONVERTIO_API = 'https://api.convertio.co/convert'
API_KEY = '3c50e707584d2cbe0139d35033b99d7c'

# Ø¥Ù†Ø´Ø§Ø¡ Ù…Ø«ÙŠÙ„ Ù„Ù„Ù…ØªØ±Ø¬Ù…
translator = Translator()

# Ù„ØªØªØ¨Ø¹ Ø¹Ø¯Ø¯ Ø§Ù„Ù…Ù„ÙØ§Øª Ø§Ù„Ù…Ø­ÙˆÙ„Ø© ÙŠÙˆÙ…ÙŠØ§Ù‹ Ù„ÙƒÙ„ Ù…Ø³ØªØ®Ø¯Ù… (user_id: (last_date, count))
user_file_usage = {}

def fix_arabic(text):
    reshaped = arabic_reshaper.reshape(text)
    return get_display(reshaped)

def translate_text_group(text_group):
    marker = "<<<SEP>>>"
    combined = marker.join(segment.strip() for segment in text_group)
    try:
        translated_combined = translator.translate(combined, src='en', dest='ar').text
        translated_combined = fix_arabic(translated_combined)
    except Exception as e:
        logger.error(f"Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ ØªØ±Ø¬Ù…Ø© Ø§Ù„Ù…Ø¬Ù…ÙˆØ¹Ø©: {e}")
        translated_combined = None

    if translated_combined:
        parts = translated_combined.split(marker)
        if len(parts) == len(text_group):
            final_parts = []
            for orig, part in zip(text_group, parts):
                leading_spaces = orig[:len(orig) - len(orig.lstrip())]
                trailing_spaces = orig[len(orig.rstrip()):]
                final_parts.append(leading_spaces + part + trailing_spaces)
            return final_parts

    result = []
    for segment in text_group:
        try:
            t = translator.translate(segment.strip(), src='en', dest='ar').text
            t = fix_arabic(t)
        except Exception as e:
            logger.error(f"Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ ØªØ±Ø¬Ù…Ø© Ø§Ù„Ø¬Ø²Ø¡: {e}")
            t = segment
        leading_spaces = segment[:len(segment) - len(segment.lstrip())]
        trailing_spaces = segment[len(segment.rstrip()):]
        result.append(leading_spaces + t + trailing_spaces)
    return result

def process_parent_texts(parent):
    new_contents = []
    group = []
    for child in parent.contents:
        if isinstance(child, NavigableString):
            group.append(str(child))
        else:
            if group:
                translated_group = translate_text_group(group)
                for text in translated_group:
                    new_contents.append(NavigableString(text))
                group = []
            new_contents.append(child)
    if group:
        translated_group = translate_text_group(group)
        for text in translated_group:
            new_contents.append(NavigableString(text))
    parent.clear()
    for item in new_contents:
        parent.append(item)

def translate_html(html_content):
    soup = BeautifulSoup(html_content, 'html.parser')
    
    head = soup.find('head')
    if head and not head.find('meta', charset=True):
        meta_tag = soup.new_tag('meta', charset='UTF-8')
        head.insert(0, meta_tag)
    
    for tag in soup.find_all():
        if tag.name in ['script', 'style']:
            continue
        if any(isinstance(child, NavigableString) and child.strip() for child in tag.contents):
            process_parent_texts(tag)
    
    return str(soup)

def build_progress_text(progress: int) -> str:
    return f"Ø¬Ø§Ø±ÙŠ Ø§Ù„ØªØ±Ø¬Ù…Ø©... {progress}%"

def translate_docx(input_path, output_path, progress_callback=None):
    """
    ØªÙØªØ­ Ø§Ù„Ø¯Ø§Ù„Ø© Ù…Ù„Ù DOCXØŒ ÙˆØªØªØ±Ø¬Ù… Ù†ØµÙˆØµ ÙƒÙ„ ÙÙ‚Ø±Ø© Ù…Ø¹ ØªØ­Ø¯ÙŠØ« Ø§Ù„ØªÙ‚Ø¯Ù… Ø¹Ø¨Ø± callback (Ø¥Ù† ÙˆØ¬Ø¯)
    Ø«Ù… ØªØ­ÙØ¸ Ø§Ù„Ù…Ù„Ù Ø§Ù„Ù…ØªØ±Ø¬Ù….
    """
    doc = Document(input_path)
    total = len(doc.paragraphs) if doc.paragraphs else 1
    for i, para in enumerate(doc.paragraphs):
        if para.text.strip():
            try:
                translated = translator.translate(para.text, src='en', dest='ar').text
                translated = fix_arabic(translated)
                para.text = translated
            except Exception as e:
                logger.error(f"Error translating DOCX paragraph: {e}")
        if progress_callback:
            progress_callback(int(((i+1) / total) * 100))
    doc.save(output_path)

def translate_pptx(input_path, output_path, progress_callback=None):
    """
    ØªÙØªØ­ Ø§Ù„Ø¯Ø§Ù„Ø© Ù…Ù„Ù PPTX ÙˆØªØ¬Ù…Ø¹ ÙƒÙ„ Ø§Ù„Ø¹Ù†Ø§ØµØ± (shapes) Ø§Ù„ØªÙŠ ØªØ­ØªÙˆÙŠ Ø¹Ù„Ù‰ Ù†ØµØŒ Ø«Ù… ØªØªØ±Ø¬Ù… ÙƒÙ„ Ù†Øµ Ù…Ø¹ ØªØ­Ø¯ÙŠØ« Ø§Ù„ØªÙ‚Ø¯Ù….
    Ø«Ù… ØªØ­ÙØ¸ Ø§Ù„Ù…Ù„Ù Ø§Ù„Ù…ØªØ±Ø¬Ù….
    """
    prs = Presentation(input_path)
    shapes_list = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shapes_list.append(shape)
    total = len(shapes_list) if shapes_list else 1
    for i, shape in enumerate(shapes_list):
        try:
            translated = translator.translate(shape.text, src='en', dest='ar').text
            translated = fix_arabic(translated)
            shape.text = translated
        except Exception as e:
            logger.error(f"Error translating PPTX shape: {e}")
        if progress_callback:
            progress_callback(int(((i+1) / total) * 100))
    prs.save(output_path)

def handle_document(update: Update, context: CallbackContext) -> None:
    # Ù…Ù†Ø¹ Ø¥Ø±Ø³Ø§Ù„ Ø£ÙƒØ«Ø± Ù…Ù† Ù…Ù„Ù ÙÙŠ Ø±Ø³Ø§Ù„Ø© ÙˆØ§Ø­Ø¯Ø©
    if update.message.media_group_id:
        update.message.reply_text("âŒ Ø§Ù„Ø±Ø¬Ø§Ø¡ Ø¥Ø±Ø³Ø§Ù„ Ù…Ù„Ù ÙˆØ§Ø­Ø¯ ÙÙ‚Ø· ÙÙŠ ÙƒÙ„ Ù…Ø±Ø©.\n Ø§Ù„Ø§ ÙˆØ³Ù ÙŠØªÙ… Ø­Ø¸Ø±ÙƒğŸ˜‚")
        return

    document = update.message.document
    if not document:
        return

    # Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø­Ø¬Ù… Ø§Ù„Ù…Ù„Ù (1 Ù…ÙŠØ¬Ø§Ø¨Ø§ÙŠØª)
    if document.file_size > 1 * 1024 * 1024:
        update.message.reply_text("âŒ Ø­Ø¬Ù… Ø§Ù„Ù…Ù„Ù Ø£ÙƒØ¨Ø± Ù…Ù† 1MB. ÙŠØ±Ø¬Ù‰ Ø¥Ø±Ø³Ø§Ù„ Ù…Ù„Ù PDF Ø£ØµØºØ±.\n Ù‚Ø³Ù… Ø¨Ø¶ØºØ· Ù…Ù„Ù ÙÙŠ Ø§Ù„Ø¨ÙˆØª Ù‡Ø°Ø§ :@i2pdfbot\n Ø«Ù… Ù‚Ù… Ø¨Ø§Ø±Ø³Ø§Ù„ Ù…Ù„Ù Ù„ÙƒÙŠ Ø§ØªØ±Ø¬Ù…Ø©")
        return

    # ØªØ­Ø¯ÙŠØ« Ø­Ø¯ Ø§Ù„Ø§Ø³ØªØ®Ø¯Ø§Ù… Ø§Ù„ÙŠÙˆÙ…ÙŠ
    user_id = update.message.from_user.id
    today_str = date.today().isoformat()
    if user_id in user_file_usage:
        last_date, count = user_file_usage[user_id]
        if last_date == today_str:
            if count >= 5:
                update.message.reply_text("ğŸš« Ù„Ù‚Ø¯ ØªØ¬Ø§ÙˆØ²Øª Ø§Ù„Ø­Ø¯ Ø§Ù„Ø£Ù‚ØµÙ‰ (5 Ù…Ù„ÙØ§Øª ÙŠÙˆÙ…ÙŠÙ‹Ø§). ÙŠØ±Ø¬Ù‰ Ø§Ù„Ù…Ø­Ø§ÙˆÙ„Ø© ØºØ¯Ù‹Ø§.")
                return
            else:
                user_file_usage[user_id] = (today_str, count + 1)
        else:
            user_file_usage[user_id] = (today_str, 1)
    else:
        user_file_usage[user_id] = (today_str, 1)

    filename_lower = document.file_name.lower()
    
    if filename_lower.endswith('.pdf'):
        # Ù…Ø¹Ø§Ù„Ø¬Ø© Ù…Ù„ÙØ§Øª PDF (ÙŠØªÙ… ØªØ­ÙˆÙŠÙ„Ù‡Ø§ Ø¥Ù„Ù‰ HTML Ø«Ù… ØªØ±Ø¬Ù…ØªÙ‡Ø§)
        file = document.get_file()
        input_filename = 'input.pdf'
        file.download(input_filename)

        try:
            with open(input_filename, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                num_pages = len(reader.pages)
            if num_pages > 5:
                update.message.reply_text("âŒ Ø§Ù„Ø­Ø¯ Ø§Ù„Ø£Ù‚ØµÙ‰ Ù‡Ùˆ 5 ØµÙØ­Ø§Øª Ø¨Ø³Ø¨Ø¨ Ø§Ù„ØªØ­Ù…ÙŠÙ„ Ø§Ù„Ø²Ø§Ø¦Ø¯.\n Ù‚Ø³Ù… Ø¨ØªÙ‚Ø³ÙŠÙ… Ù…Ù„Ù ÙÙŠ Ø§Ù„Ø¨ÙˆØª Ù‡Ø°Ø§ :@i2pdfbot\n Ø«Ù… Ù‚Ù… Ø¨Ø§Ø±Ø³Ø§Ù„ Ù…Ù„Ù Ù„ÙƒÙŠ Ø§ØªØ±Ø¬Ù…Ø©")
                os.remove(input_filename)
                return
        except Exception as e:
            logger.error(f"Error reading PDF: {e}")
            update.message.reply_text("Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ Ù‚Ø±Ø§Ø¡Ø© Ø§Ù„Ù…Ù„Ù.")
            os.remove(input_filename)
            return

        with open(input_filename, 'rb') as f:
            file_data = f.read()
        encoded_file = base64.b64encode(file_data).decode('utf-8')

        payload = {
            "apikey": API_KEY,
            "input": "base64",
            "file": encoded_file,
            "filename": document.file_name,
            "outputformat": "html"
        }

        try:
            response = requests.post(CONVERTIO_API, json=payload)
            response.raise_for_status()
        except Exception as e:
            logger.error(f"Error during conversion initiation: {e}")
            update.message.reply_text("Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ Ø¨Ø¯Ø¡ Ø¹Ù…Ù„ÙŠØ© Ø§Ù„ØªØ­ÙˆÙŠÙ„.")
            os.remove(input_filename)
            return

        result = response.json()
        if result.get('code') != 200:
            error_msg = result.get('error', 'Ø®Ø·Ø£ ØºÙŠØ± Ù…Ø¹Ø±ÙˆÙ.')
            update.message.reply_text(f"Ø®Ø·Ø£ ÙÙŠ API Ø§Ù„ØªØ­ÙˆÙŠÙ„: {error_msg}")
            os.remove(input_filename)
            return

        conversion_id = result['data']['id']
        status_url = f"{CONVERTIO_API}/{conversion_id}/status"
        start_time = time.time()
        max_wait_time = 60
        progress_message = update.message.reply_text("Ø¬Ø§Ø±ÙŠ Ø§Ù„ØªØ±Ø¬Ù…Ø© Ø§Ù„Ø«Ù†Ø§Ø¦ÙŠØ©... 0%")
        
        while True:
            time.sleep(2)
            elapsed = time.time() - start_time
            progress = min(int((elapsed / max_wait_time) * 100), 100)
            try:
                status_resp = requests.get(status_url)
                status_data = status_resp.json()
            except Exception as e:
                logger.error(f"Error checking conversion status: {e}")
                update.message.reply_text("Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ Ø§Ù„ØªØ­Ù‚Ù‚ Ù…Ù† Ø­Ø§Ù„Ø© Ø§Ù„ØªØ­ÙˆÙŠÙ„.")
                os.remove(input_filename)
                return
            step = status_data.get('data', {}).get('step')
            try:
                context.bot.edit_message_text(chat_id=update.message.chat_id,
                                              message_id=progress_message.message_id,
                                              text=f"Ø¬Ø§Ø±ÙŠ Ø§Ù„ØªØ±Ø¬Ù…Ø© Ø§Ù„Ø«Ù†Ø§Ø¦ÙŠØ©... {progress}%")
            except Exception as e:
                logger.error(f"Error editing progress message: {e}")
            if step == 'finish':
                try:
                    context.bot.edit_message_text(chat_id=update.message.chat_id,
                                                  message_id=progress_message.message_id,
                                                  text="Ø¬Ø§Ø±ÙŠ Ø§Ù„ØªØ±Ø¬Ù…Ø© Ø§Ù„Ø«Ù†Ø§Ø¦ÙŠØ©... 100%")
                except Exception as e:
                    logger.error(f"Error finalizing progress message: {e}")
                break
            if step == 'error':
                update.message.reply_text("Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ Ø§Ù„ØªØ­ÙˆÙŠÙ„.")
                os.remove(input_filename)
                return

        try:
            download_url = status_data['data']['output']['url']
            download_resp = requests.get(download_url)
            download_resp.raise_for_status()
        except Exception as e:
            logger.error(f"Error downloading converted file: {e}")
            update.message.reply_text("Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ ØªØ­Ù…ÙŠÙ„ Ø§Ù„Ù…Ù„Ù Ø§Ù„Ù…Ø­ÙˆÙ„.")
            os.remove(input_filename)
            return

        output_filename = 'output.html'
        with open(output_filename, 'wb') as f:
            f.write(download_resp.content)

        try:
            with open(output_filename, 'r', encoding='utf-8') as f:
                html_content = f.read()
        except Exception as e:
            logger.error(f"Error reading converted HTML: {e}")
            update.message.reply_text("Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ Ù‚Ø±Ø§Ø¡Ø© Ø§Ù„Ù…Ù„Ù Ø§Ù„Ù…Ø­ÙˆÙ„.")
            os.remove(input_filename)
            os.remove(output_filename)
            return

        # ØªØ±Ø¬Ù…Ø© Ù…Ø­ØªÙˆÙ‰ HTML Ø§Ù„Ù†Ø§ØªØ¬
        translated_html = translate_html(html_content)
        base_name = os.path.splitext(document.file_name)[0]
        translated_file_path = f"{base_name}.html"
        with open(translated_file_path, 'w', encoding='utf-8') as f:
            f.write(translated_html)

        update.message.reply_document(document=open(translated_file_path, 'rb'),
                                      caption="âœ… ØªÙ… ØªØ±Ø¬Ù…Ø© Ø§Ù„Ù…Ù„Ù Ø¨Ù†Ø¬Ø§Ø­!\n Ø§Ø°Ø§ Ù„Ù… ÙŠØ¹Ø¬Ø¨Ùƒ ØªØµÙ…ÙŠÙ…  Ù…Ù„Ù ÙÙŠ Ø§Ø¹Ù„Ø§ ÙŠÙ…ÙƒÙ†Ùƒ Ø§Ø³ØªØ®Ø¯Ø§Ù… Ù‡Ø°Ø§ Ù…Ù„Ù\n @ta_ja199 Ù„Ø§Ø³ØªÙØ³Ø§Ø±")
        context.bot.delete_message(chat_id=update.message.chat_id,
                                   message_id=progress_message.message_id)

        os.remove(input_filename)
        os.remove(output_filename)
        os.remove(translated_file_path)

    elif filename_lower.endswith('.docx'):
        # Ù…Ø¹Ø§Ù„Ø¬Ø© Ù…Ù„ÙØ§Øª DOCX: ØªØ±Ø¬Ù…Ø© Ø§Ù„Ù…Ù„Ù Ù…Ø¨Ø§Ø´Ø±Ø© Ù…Ø¹ ØªØ­Ø¯ÙŠØ« Ø±Ø³Ø§Ù„Ø© Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø±
        input_filename = 'input.docx'
        document.get_file().download(input_filename)
        base_name = os.path.splitext(document.file_name)[0]
        output_filename = f"{base_name}.docx"

        progress_message = update.message.reply_text("Ø¬Ø§Ø±ÙŠ Ø§Ù„ØªØ±Ø¬Ù…Ø© Ø§Ù„Ø«Ù†Ø§Ø¦ÙŠØ©... 0%")
        def progress_callback(progress):
            try:
                context.bot.edit_message_text(chat_id=update.message.chat_id,
                                              message_id=progress_message.message_id,
                                              text=f"Ø¬Ø§Ø±ÙŠ Ø§Ù„ØªØ±Ø¬Ù…Ø©... {progress}%")
            except Exception as e:
                logger.error(f"Error updating progress message for DOCX: {e}")

        try:
            translate_docx(input_filename, output_filename, progress_callback)
        except Exception as e:
            logger.error(f"Error processing DOCX: {e}")
            update.message.reply_text("Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ ØªØ±Ø¬Ù…Ø© Ù…Ù„Ù DOCX.")
            os.remove(input_filename)
            return

        try:
            context.bot.delete_message(chat_id=update.message.chat_id,
                                       message_id=progress_message.message_id)
        except Exception as e:
            logger.error(f"Error deleting progress message for DOCX: {e}")

        update.message.reply_document(document=open(output_filename, 'rb'),
                                      caption="âœ… ØªÙ… ØªØ±Ø¬Ù…Ø© Ø§Ù„Ù…Ù„Ù Ø¨Ù†Ø¬Ø§Ø­!\n ÙŠÙ…ÙƒÙ†Ùƒ Ø§Ø³ØªØ¹Ù…Ø§Ù„ Ù‡Ø°Ø§ Ø§Ù„Ø¨ÙˆØª ÙÙŠ ØªØ­ÙˆÙŠÙ„Ù‡ Ù„pdf :@i2pdfbot")
        os.remove(input_filename)
        os.remove(output_filename)

    elif filename_lower.endswith('.pptx'):
        # Ù…Ø¹Ø§Ù„Ø¬Ø© Ù…Ù„ÙØ§Øª PPTX: ØªØ±Ø¬Ù…Ø© Ø§Ù„Ù…Ù„Ù Ù…Ø¨Ø§Ø´Ø±Ø© Ù…Ø¹ ØªØ­Ø¯ÙŠØ« Ø±Ø³Ø§Ù„Ø© Ø§Ù„Ø§Ù†ØªØ¸Ø§Ø±
        input_filename = 'input.pptx'
        document.get_file().download(input_filename)
        base_name = os.path.splitext(document.file_name)[0]
        output_filename = f"{base_name}.pptx"

        progress_message = update.message.reply_text("Ø¬Ø§Ø±ÙŠ Ø§Ù„ØªØ±Ø¬Ù…Ø©... 0%")
        def progress_callback(progress):
            try:
                context.bot.edit_message_text(chat_id=update.message.chat_id,
                                              message_id=progress_message.message_id,
                                              text=f"Ø¬Ø§Ø±ÙŠ Ø§Ù„ØªØ±Ø¬Ù…Ø©... {progress}%")
            except Exception as e:
                logger.error(f"Error updating progress message for PPTX: {e}")

        try:
            translate_pptx(input_filename, output_filename, progress_callback)
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            update.message.reply_text("Ø­Ø¯Ø« Ø®Ø·Ø£ Ø£Ø«Ù†Ø§Ø¡ ØªØ±Ø¬Ù…Ø© Ù…Ù„Ù PPTX.")
            os.remove(input_filename)
            return

        try:
            context.bot.delete_message(chat_id=update.message.chat_id,
                                       message_id=progress_message.message_id)
        except Exception as e:
            logger.error(f"Error deleting progress message for PPTX: {e}")

        update.message.reply_document(document=open(output_filename, 'rb'),
                                      caption="âœ… ØªÙ… ØªØ±Ø¬Ù…Ø© Ø§Ù„Ù…Ù„Ù Ø¨Ù†Ø¬Ø§Ø­!\n ÙŠÙ…ÙƒÙ†Ùƒ Ø§Ø³ØªØ¹Ù…Ø§Ù„ Ù‡Ø°Ø§ Ø§Ù„Ø¨ÙˆØª ÙÙŠ ØªØ­ÙˆÙŠÙ„Ù‡ Ù„pdf :@i2pdfbot")
        os.remove(input_filename)
        os.remove(output_filename)

    else:
        update.message.reply_text("ÙŠØ±Ø¬Ù‰ Ø¥Ø±Ø³Ø§Ù„ Ù…Ù„Ù Ø¨ØµÙŠØºØ© PDF, DOCX, Ø£Ùˆ PPTX ÙÙ‚Ø·.")

def help(update: Update, context: CallbackContext) -> None:
    update.message.reply_text(
        "Ù…Ø±Ø­Ø¨Ø§Ù‹ Ø¨Ùƒ Ù‚Ø³Ù… Ù…Ø³Ø§Ø¹Ø¯Ø©! ÙŠØ±Ø¬Ù‰ Ø¥Ø±Ø³Ø§Ù„ Ù…Ù„Ù Ù…Ù† Ø§Ù„Ø£Ù†ÙˆØ§Ø¹ Ø§Ù„ØªØ§Ù„ÙŠØ©:\n"
        "ÙŠÙ…ÙƒÙ† ØªØ±Ø¬Ù…Ø© Ù…Ù„ÙØ§Øª pdf Ù„ÙƒÙ† Ø¨Ù‚ÙŠÙˆØ¯ Ù„ÙƒÙŠ Ù„Ø§ ÙŠØªÙˆÙ‚Ù Ø§Ù„Ø¨ÙˆØª ØªÙ… ØªÙ‚ÙŠÙŠØ¯ Ù…Ù„ÙØ§Øª Ø¨\n"
        "1BM ÙƒØ­Ø¯ Ø§Ù‚ØµÙ‰ \nÙˆØ­Ø¯ 5 ØµÙØ­Ø§Øª Ùˆ Ø®Ù…Ø³ Ù…Ù„ÙØ§Øª ÙÙŠ Ø§Ù„ÙŠÙˆÙ…\n"
        "Ø§Ù„ØµÙŠØº Ø§Ù„ØªÙŠ ÙŠÙ…ÙƒÙ† ØªØ±Ø¬Ù…ØªÙ‡Ø§ Ù‡ÙŠ:pdf,docx,pptx\n"
        "Ù„Ø§Ø³ØªÙØ³Ø§Ø± @ta_ja199"
    )

def main() -> None:
    updater = Updater(TELEGRAM_TOKEN, use_context=True)
    dp = updater.dispatcher

    dp.add_handler(CommandHandler("help", help))
    dp.add_handler(MessageHandler(Filters.document, handle_document))

    updater.start_polling()
    logger.info("Ø§Ù„Ø¨ÙˆØª ÙŠØ¹Ù…Ù„ Ø§Ù„Ø¢Ù†...")
    updater.idle()

if __name__ == '__main__':
    main()
